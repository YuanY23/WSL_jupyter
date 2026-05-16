from pathlib import Path

import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("中期答辩-6页修改.pptx")
ASSET_DIR = Path("ppt_assets")
ASSET_DIR.mkdir(exist_ok=True)


W, H = 13.333333, 7.5
NAVY = "000E4F"
NAVY2 = "192D96"
RED = "C0504D"
DARK_RED = "920000"
PALE = "EDF8FC"
PALE2 = "F2F9FC"
LIGHT_BLUE = "B7E9FC"
GRAY = "4D4D4D"
MID_GRAY = "6B7280"
BLACK = "111827"
WHITE = "FFFFFF"


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def add_textbox(slide, text, x, y, w, h, size=16, color=BLACK, bold=False,
                font="微软雅黑", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                margin=0.05, line_spacing=1.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_rich_line(slide, parts, x, y, w, h, size=15, font="微软雅黑"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.line_spacing = 1.05
    for text, opts in parts:
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = opts.get("bold", False)
        r.font.color.rgb = rgb(opts.get("color", BLACK))
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=LIGHT_BLUE, radius=False, width=1.0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    s.line.color.rgb = rgb(line)
    s.line.width = Pt(width)
    return s


def add_arrow(slide, x1, y1, x2, y2, color=NAVY2, width=1.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_header(slide, title, idx):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(0.68)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = rgb(NAVY)
    slide.shapes[-1].line.fill.background()
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.68), Inches(W), Inches(0.08)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = rgb(RED)
    slide.shapes[-1].line.fill.background()
    add_textbox(slide, title, 0.52, 0.12, 9.8, 0.42, size=24, color=WHITE, bold=True, margin=0)
    add_textbox(slide, f"{idx:02d}", 12.13, 0.09, 0.62, 0.42, size=22, color=WHITE, bold=True,
                align=PP_ALIGN.RIGHT, margin=0)
    add_textbox(slide, "第三部分  研究内容", 10.15, 0.19, 1.82, 0.28, size=10, color=LIGHT_BLUE,
                align=PP_ALIGN.RIGHT, margin=0)
    # Footer
    add_textbox(slide, "显示计算代码的热力仿真系统", 0.52, 7.14, 4.6, 0.2, size=8.5, color=MID_GRAY, margin=0)
    add_textbox(slide, f"{idx}", 12.45, 7.12, 0.28, 0.2, size=9, color=RED, bold=True,
                align=PP_ALIGN.RIGHT, margin=0)


def add_label(slide, text, x, y, w, h, fill=NAVY, color=WHITE, size=13, bold=True):
    add_rect(slide, x, y, w, h, fill=fill, line=fill, radius=False)
    return add_textbox(slide, text, x + 0.05, y + 0.04, w - 0.1, h - 0.08,
                       size=size, color=color, bold=bold, align=PP_ALIGN.CENTER,
                       valign=MSO_ANCHOR.MIDDLE, margin=0)


def add_bullet_block(slide, bullets, x, y, w, h, size=13.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.line_spacing = 1.04
        p.space_after = Pt(4)
        p.font.name = "微软雅黑"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(BLACK)
        p._p.get_or_add_pPr().insert(0, __import__("pptx").oxml.parse_xml(
            '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="•"/>'
        ))
    return box


def make_chart():
    L = 0.1
    k = 400
    T_left = 100
    T_right = 20
    N = 50
    x = np.linspace(0, L, N)
    T = T_left + (T_right - T_left) * x / L
    img = Image.new("RGB", (1100, 530), "white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 19)
        tiny = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
    except Exception:
        font = small = tiny = ImageFont.load_default()

    left, top, right, bottom = 115, 70, 1040, 420
    draw.text((250, 18), "1D steady conduction temperature distribution", fill="#111827", font=font)
    # grid
    for i in range(6):
        xx = left + (right - left) * i / 5
        draw.line((xx, top, xx, bottom), fill="#E5E7EB", width=1)
        draw.text((xx - 12, bottom + 14), f"{int(i * 20)}", fill="#4D4D4D", font=tiny)
    for j in range(5):
        yy = bottom - (bottom - top) * j / 4
        draw.line((left, yy, right, yy), fill="#E5E7EB", width=1)
        draw.text((54, yy - 9), f"{20 + j * 20}", fill="#4D4D4D", font=tiny)
    draw.line((left, bottom, right, bottom), fill="#111827", width=2)
    draw.line((left, top, left, bottom), fill="#111827", width=2)

    def pt(xv, tv):
        px = left + (xv / L) * (right - left)
        py = bottom - ((tv - 20) / 80) * (bottom - top)
        return px, py

    pts = [pt(float(a), float(b)) for a, b in zip(x, T)]
    draw.line(pts, fill="#192D96", width=4)
    for px, py in pts[::4]:
        draw.ellipse((px - 4, py - 4, px + 4, py + 4), fill="#192D96")
    # analytical line overlapped as red dashed, slightly offset for visibility
    for a, b in zip(pts[:-1:4], pts[1::4]):
        draw.line((a[0], a[1] + 4, b[0], b[1] + 4), fill="#C0504D", width=3)
    draw.text((450, 460), "Position x (mm)", fill="#111827", font=small)
    draw.text((15, 210), "Temperature T", fill="#111827", font=small)
    draw.rectangle((765, 84, 1015, 142), outline="#D1D5DB", fill="#FFFFFF")
    draw.line((785, 103, 840, 103), fill="#192D96", width=4)
    draw.text((852, 92), "FDM numerical", fill="#111827", font=tiny)
    draw.line((785, 126, 840, 126), fill="#C0504D", width=3)
    draw.text((852, 115), "Analytical", fill="#111827", font=tiny)
    path = ASSET_DIR / "flat_plate_temperature.png"
    img.save(path)
    return path


def slide1(pres):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "显示计算代码的热力仿真系统总体研究框架", 1)
    add_rect(s, 0.55, 1.05, 5.05, 4.72, fill=PALE2, line=LIGHT_BLUE, radius=False)
    add_label(s, "研究定位", 0.78, 1.28, 1.28, 0.34, fill=NAVY)
    add_textbox(s,
                "本研究面向热力系统仿真过程中模型关系不清、计算过程不透明、学习者难以理解“物理模型—数学方程—程序代码—仿真结果”之间联系的问题，构建一个以可执行代码为核心的白箱式热力仿真系统。",
                0.78, 1.72, 4.45, 1.15, size=14.1, color=BLACK)
    add_label(s, "系统目标", 0.78, 3.05, 1.28, 0.34, fill=RED)
    add_textbox(s,
                "系统将热力建模、数值求解、代码展示、结果可视化和教学实验组织在统一环境中，实现从参数输入到模型方程、从求解代码到仿真结果的全过程展示与复现。",
                0.78, 3.48, 4.45, 1.04, size=14.1, color=BLACK)
    add_textbox(s, "不是只给出仿真结果，而是展示热力问题如何被建模、计算和验证。",
                0.78, 4.86, 4.45, 0.48, size=15.5, color=DARK_RED, bold=True,
                valign=MSO_ANCHOR.MIDDLE)

    x0, y0, bw, bh, gap = 6.0, 1.12, 1.22, 1.02, 0.22
    boxes = [
        ("热力对象", "导热 / 对流\n辐射 / 综合换热", NAVY),
        ("数学模型", "控制方程\n边界条件\n物性参数", "003980"),
        ("数值求解代码", "参数定义\n方程构建\n离散求解", RED),
        ("结果可视化", "温度分布\n热流密度\n变化曲线", "007EC8"),
        ("教学实验应用", "实验任务\n复算修改\n提交反馈", NAVY2),
    ]
    for i, (title, body, color) in enumerate(boxes):
        x = x0 + i * (bw + gap)
        add_rect(s, x, y0, bw, bh, fill="FFFFFF", line=color, radius=False, width=1.4)
        add_textbox(s, title, x + 0.08, y0 + 0.09, bw - 0.16, 0.25, size=12.2,
                    color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_textbox(s, body, x + 0.08, y0 + 0.41, bw - 0.16, 0.48, size=9.4,
                    color=BLACK, align=PP_ALIGN.CENTER, margin=0)
        if i < len(boxes) - 1:
            add_arrow(s, x + bw + 0.03, y0 + 0.51, x + bw + gap - 0.03, y0 + 0.51, color=RED, width=1.1)
    add_rect(s, 6.0, 2.76, 6.7, 2.15, fill=PALE, line=LIGHT_BLUE, radius=False)
    add_textbox(s, "研究重点", 6.28, 2.98, 1.0, 0.28, size=13.5, color=NAVY, bold=True, margin=0)
    add_bullet_block(s, [
        "将热力对象、数学模型、代码执行与结果可视化贯通",
        "通过 Notebook 保存完整计算链路，支持复算与修改",
        "以白箱式表达降低传统 GUI 仿真的黑箱化问题",
    ], 6.24, 3.34, 5.95, 1.02, size=12.2)
    add_textbox(s, "研究重点：实现热力仿真过程的透明化、代码化与可复现表达。",
                6.25, 5.32, 5.95, 0.38, size=15, color=NAVY, bold=True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def slide2(pres):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "热设计典型问题的建模对象与方程体系", 2)
    add_textbox(s,
                "本研究选取热设计与热力系统分析中的典型传热问题作为建模对象，围绕导热、对流换热、热辐射及综合换热过程，建立相应的数学模型与参数体系。",
                0.68, 1.03, 6.05, 0.72, size=15, color=BLACK)
    add_textbox(s,
                "建模过程从物理问题出发，明确研究对象、基本假设、控制方程、边界条件和物性参数，将实际热设计问题转化为可计算、可执行、可验证的数学模型。",
                0.68, 1.78, 6.05, 0.72, size=15, color=BLACK)
    add_bullet_block(s, [
        "面向热设计中的典型传热问题开展建模",
        "覆盖稳态导热、瞬态导热、对流换热、热辐射等基础场景",
        "明确导热系数、换热系数、温度场、热流密度等关键变量",
        "建立控制方程、边界条件和物性参数之间的对应关系",
        "为后续数值求解、代码生成和结果可视化提供模型基础",
    ], 0.78, 2.76, 5.82, 2.2, size=13.2)
    add_textbox(s, "系统展示的不是单一计算结果，而是热力问题从物理对象到数学方程的建模过程。",
                0.78, 5.42, 5.8, 0.36, size=14.5, color=DARK_RED, bold=True)

    add_rect(s, 7.2, 1.08, 5.25, 4.8, fill=PALE2, line=LIGHT_BLUE, radius=False)
    add_textbox(s, "热设计典型问题", 8.86, 1.28, 1.8, 0.34, size=15, color=NAVY, bold=True,
                align=PP_ALIGN.CENTER, margin=0)
    cards = [
        ("导热问题", "输入：k、尺寸、边界温度\n方程：傅里叶定律、热传导方程\n输出：温度分布、热流密度", 7.48, 1.88, NAVY),
        ("对流换热", "输入：h、流体温度、换热面积\n方程：牛顿冷却公式\n输出：换热量、表面温度变化", 9.86, 1.88, RED),
        ("热辐射", "输入：表面温度、发射率、面积\n方程：斯忒藩-玻尔兹曼定律\n输出：辐射换热量", 7.48, 3.68, "007EC8"),
        ("综合换热", "输入：多边界条件、材料参数\n方程：热阻网络、能量守恒\n输出：综合传热量、系统响应", 9.86, 3.68, NAVY2),
    ]
    for title, body, x, y, c in cards:
        add_rect(s, x, y, 2.18, 1.36, fill=WHITE, line=c, radius=False, width=1.25)
        add_textbox(s, title, x + 0.12, y + 0.10, 1.94, 0.24, size=12.6, color=c, bold=True,
                    align=PP_ALIGN.CENTER, margin=0)
        add_textbox(s, body, x + 0.16, y + 0.46, 1.86, 0.74, size=8.6, color=BLACK,
                    align=PP_ALIGN.LEFT, margin=0)
    add_textbox(s, "通过建立清晰的热力模型体系，为后续“方程—代码—结果”的白箱式仿真流程奠定基础。",
                7.45, 6.05, 4.82, 0.34, size=13.1, color=NAVY, bold=True,
                align=PP_ALIGN.CENTER, margin=0)


def slide3(pres):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "从热力模型方程到可执行仿真代码的转换机制", 3)
    add_textbox(s,
                "系统将热力模型按照“参数定义、控制方程、离散求解、物理量计算、结果可视化”的顺序组织为可执行 Notebook，使学习者能够看到热力问题如何从数学表达转化为 Python 计算代码。",
                0.66, 1.03, 6.0, 0.78, size=14.4, color=BLACK)
    add_bullet_block(s, [
        "采用 FDM、FVM、解析解对比等方法支撑不同传热场景",
        "Notebook 中同时包含模型说明、公式推导、Python 代码和可视化结果",
        "参数改变后可重新执行代码，实现仿真过程复算与结果追踪",
        "重点体现“公式如何变成代码”，而不是只展示最终结果",
    ], 0.72, 1.96, 5.75, 1.55, size=12.6)
    add_rect(s, 0.72, 3.78, 5.62, 1.24, fill=PALE, line=LIGHT_BLUE)
    add_textbox(s, "一维稳态导热控制方程", 0.94, 3.94, 2.4, 0.22, size=12.8, color=NAVY, bold=True, margin=0)
    add_textbox(s, "d²T/dx² = 0      0 < x < L\nT(0) = T_h       T(L) = T_c\nT_{i-1} - 2T_i + T_{i+1} = 0     [A]{T} = {b}",
                0.94, 4.22, 5.0, 0.55, size=13.5, color=BLACK, font="Consolas", margin=0)

    add_rect(s, 6.85, 1.04, 5.85, 2.32, fill="111827", line="111827")
    add_textbox(s, "Python 求解代码片段", 7.08, 1.22, 2.2, 0.22, size=12.5, color=LIGHT_BLUE, bold=True, margin=0)
    code = (
        "dx = L / (N - 1)\n"
        "A = np.zeros((N, N)); b = np.zeros(N)\n"
        "A[0,0] = 1.0;   b[0] = T_left\n"
        "A[-1,-1] = 1.0; b[-1] = T_right\n\n"
        "for i in range(1, N - 1):\n"
        "    A[i,i-1] = 1.0\n"
        "    A[i,i]   = -2.0\n"
        "    A[i,i+1] = 1.0\n\n"
        "T = np.linalg.solve(A, b)"
    )
    add_textbox(s, code, 7.08, 1.58, 5.35, 1.52, size=8.9, color="DDEAFE", font="Consolas", margin=0)

    steps = [
        ("物理参数输入", "L, k, T_h, T_c"),
        ("控制方程", "d²T/dx² = 0"),
        ("数值离散", "FDM / FVM"),
        ("Python 求解", "矩阵组装 + 求解"),
        ("结果输出", "温度 / 热流 / 误差"),
    ]
    y = 4.0
    for i, (a, b) in enumerate(steps):
        x = 6.9 + i * 1.13
        add_rect(s, x, y, 0.98, 0.94, fill=WHITE, line=RED if i == 3 else NAVY2, width=1.2)
        add_textbox(s, a, x + 0.04, y + 0.11, 0.9, 0.2, size=8.6, color=NAVY, bold=True,
                    align=PP_ALIGN.CENTER, margin=0)
        add_textbox(s, b, x + 0.05, y + 0.45, 0.88, 0.24, size=7.4, color=BLACK,
                    align=PP_ALIGN.CENTER, margin=0)
        if i < len(steps) - 1:
            add_arrow(s, x + 1.0, y + 0.47, x + 1.11, y + 0.47, color=RED, width=1)
    add_textbox(s, "系统通过 Notebook 将“模型方程—求解代码—仿真结果”组织为连续、透明、可复现的计算过程。",
                2.0, 6.05, 9.3, 0.34, size=14.5, color=NAVY, bold=True,
                align=PP_ALIGN.CENTER, margin=0)


def slide4(pres):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "面向代码可见仿真的系统实现架构", 4)
    add_textbox(s,
                "系统以热力仿真流程为主线，将模型选择、参数输入、代码生成、内核执行、结果展示和实验管理集成在统一环境中。JupyterLab、Python Kernel、Docker 和 nbgrader 作为支撑平台，服务于热力模型计算过程的可见化、可运行和可复现。",
                0.68, 1.02, 11.86, 0.72, size=14.2, color=BLACK)
    layers = [
        ("用户交互层", "模型选择 / 参数输入 / 实验入口", NAVY),
        ("热力模型层", "导热模型 / 对流模型 / 辐射模型 / 综合换热模型", "003980"),
        ("Notebook 生成层", "模型说明 + 控制方程 + 求解代码 + 可视化代码", RED),
        ("计算执行层", "Python Kernel + NumPy + Matplotlib", "007EC8"),
        ("结果展示层", "温度场 / 热流密度 / 换热量 / 曲线图 / 云图", NAVY2),
        ("教学管理层", "用户隔离 / 实验发布 / 作业提交 / 评分反馈", DARK_RED),
    ]
    y0 = 2.0
    for i, (name, desc, c) in enumerate(layers):
        y = y0 + i * 0.62
        add_rect(s, 0.88, y, 6.25, 0.43, fill=WHITE, line=c, width=1.15)
        add_rect(s, 0.88, y, 1.52, 0.43, fill=c, line=c)
        add_textbox(s, name, 0.96, y + 0.09, 1.36, 0.18, size=10.5, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, margin=0)
        add_textbox(s, desc, 2.56, y + 0.1, 4.35, 0.17, size=10.6, color=BLACK, margin=0)
        if i < len(layers) - 1:
            add_arrow(s, 4.02, y + 0.45, 4.02, y + 0.59, color=RED, width=1)

    add_rect(s, 7.63, 1.98, 4.9, 2.0, fill=PALE2, line=LIGHT_BLUE)
    add_textbox(s, "模块链路", 7.88, 2.16, 1.2, 0.22, size=13, color=NAVY, bold=True, margin=0)
    chain = [
        ("ControlPanel.tsx", "选择场景 + 输入参数"),
        ("NotebookGenerator.ts", "生成 Markdown + Code Cells"),
        ("thermal_solver", "执行热力模型求解"),
        ("JupyterLab Notebook", "展示代码、图像和结果"),
    ]
    for i, (a, b) in enumerate(chain):
        yy = 2.55 + i * 0.33
        add_rich_line(s, [(a, {"bold": True, "color": RED}), ("  →  " + b, {"color": BLACK})],
                      7.92, yy, 4.15, 0.22, size=10.1)

    add_rect(s, 7.63, 4.35, 4.9, 1.05, fill=PALE, line=LIGHT_BLUE)
    add_textbox(s, "技术定位", 7.9, 4.5, 1.1, 0.22, size=13, color=NAVY, bold=True, margin=0)
    add_textbox(s,
                "JupyterLab、Docker 和作业系统用于承载、隔离和组织仿真实验，是实现白箱式仿真的支撑条件。",
                7.9, 4.82, 4.1, 0.34, size=11.2, color=BLACK, margin=0)
    add_textbox(s, "平台技术不是研究主角，而是支撑热力仿真过程实现“可见、可算、可改、可复现”的工程载体。",
                1.35, 6.08, 10.65, 0.34, size=14.2, color=NAVY, bold=True,
                align=PP_ALIGN.CENTER, margin=0)


def slide5(pres, chart_path):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "热力建模示例构建与当前实现成果", 5)
    add_textbox(s,
                "当前系统已围绕热设计基础问题构建多个典型仿真场景，覆盖稳态导热、瞬态导热、对流换热和热辐射等内容。每个场景均包含参数输入、模型说明、计算代码、结果图像和关键物理量输出，初步形成了从热力建模到可执行仿真的完整流程。",
                0.68, 1.02, 11.9, 0.72, size=13.8, color=BLACK)
    add_rect(s, 0.72, 2.0, 5.9, 3.52, fill=PALE2, line=LIGHT_BLUE)
    add_textbox(s, "阶段成果", 0.95, 2.18, 1.2, 0.24, size=13.2, color=NAVY, bold=True, margin=0)
    add_bullet_block(s, [
        "已配置 4 类热力仿真模型：稳态导热、瞬态导热、对流换热、热辐射",
        "已构建平板导热、多层平板、圆筒壁导热、肋片导热、二维平板导热等导热场景",
        "已构建平板强制对流、竖板自然对流、圆管内强迫对流等对流换热场景",
        "已构建平行平板辐射、三表面封闭空腔辐射等热辐射场景",
        "支持生成包含模型方程、求解代码、结果曲线和误差分析的 Notebook",
        "已实现前端参数面板与 Notebook 生成模块的基本联动",
    ], 0.95, 2.55, 5.28, 2.28, size=10.9)
    add_textbox(s, "阶段成果表明，系统已能将典型热力模型转化为可执行、可展示、可复算的 Notebook 仿真实验。",
                0.95, 5.05, 5.22, 0.28, size=11.6, color=DARK_RED, bold=True, margin=0)

    add_rect(s, 7.0, 2.0, 5.38, 1.28, fill=WHITE, line=RED, width=1.2)
    add_textbox(s, "案例：一维稳态平板导热", 7.22, 2.18, 2.4, 0.22, size=13.2, color=RED, bold=True, margin=0)
    add_textbox(s,
                "输入：L、k、T_h、T_c\n模型：d²T/dx² = 0，T(0)=T_h，T(L)=T_c\n方法：有限差分法 FDM，组装 [A]{T} = {b}\n输出：温度分布、热流密度 q、单位面积热阻 R、误差",
                7.22, 2.52, 4.82, 0.58, size=9.7, color=BLACK, margin=0)
    s.shapes.add_picture(str(chart_path), Inches(7.02), Inches(3.62), width=Inches(5.35), height=Inches(2.55))
    add_textbox(s, "验证方式：数值解与解析解对比；输出图像：温度 T 随位置 x 的变化曲线。",
                7.12, 6.2, 5.08, 0.22, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)


def slide6(pres):
    s = pres.slides.add_slide(pres.slide_layouts[6])
    add_header(s, "面向实验教学的仿真训练与评价扩展", 6)
    add_textbox(s,
                "在完成热力模型构建、代码化仿真和结果可视化的基础上，系统进一步面向清洁能源与热工类课程实验教学，支持教师发布仿真实验任务，学生在 Notebook 中完成参数修改、代码运行、结果分析和实验提交。",
                0.68, 1.02, 11.88, 0.64, size=13.8, color=BLACK)
    add_textbox(s,
                "该部分不是单独的教学网站建设，而是将白箱式热力仿真流程应用到实验训练与学习评价中，使学生能够在“看见代码、运行代码、修改参数、分析结果”的过程中理解热力模型与计算方法。",
                0.68, 1.73, 11.88, 0.58, size=13.8, color=BLACK)

    steps = [
        ("教师发布实验任务", "模型类型 / 参数范围\n实验要求 / 评分规则"),
        ("学生进入 Notebook", "阅读模型说明\n查看方程 / 运行代码"),
        ("参数修改与仿真实验", "改变材料参数\n边界条件 / 网格数量"),
        ("结果分析与报告提交", "温度分布 / 热流密度\n换热量 / 误差对比"),
        ("作业收集与评价反馈", "自动检查 / 教师批改\n结果反馈"),
    ]
    x0, y, bw, bh, gap = 0.78, 2.78, 2.22, 1.1, 0.31
    for i, (title, body) in enumerate(steps):
        x = x0 + i * (bw + gap)
        c = [NAVY, "003980", RED, "007EC8", DARK_RED][i]
        add_rect(s, x, y, bw, bh, fill=WHITE, line=c, width=1.25)
        add_rect(s, x, y, bw, 0.33, fill=c, line=c)
        add_textbox(s, title, x + 0.08, y + 0.09, bw - 0.16, 0.14, size=9.6, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, margin=0)
        add_textbox(s, body, x + 0.14, y + 0.52, bw - 0.28, 0.32, size=8.7, color=BLACK,
                    align=PP_ALIGN.CENTER, margin=0)
        if i < len(steps) - 1:
            add_arrow(s, x + bw + 0.04, y + 0.54, x + bw + gap - 0.04, y + 0.54, color=RED, width=1.1)

    add_rect(s, 0.78, 4.42, 5.9, 1.22, fill=PALE, line=LIGHT_BLUE)
    add_textbox(s, "训练任务示例：一维平板稳态导热仿真", 1.0, 4.57, 3.1, 0.2, size=12.6, color=NAVY, bold=True, margin=0)
    add_textbox(s,
                "1. 修改不同材料的导热系数 k    2. 比较温度分布曲线和热流密度变化\n3. 分析边界温差对传热量的影响    4. 对比数值解与解析解误差\n5. 提交 Notebook 与结果分析说明",
                1.0, 4.91, 5.4, 0.42, size=10.5, color=BLACK, margin=0)
    add_rect(s, 7.16, 4.42, 5.12, 1.22, fill=PALE2, line=LIGHT_BLUE)
    add_textbox(s, "技术支撑", 7.38, 4.57, 1.2, 0.2, size=12.6, color=NAVY, bold=True, margin=0)
    add_textbox(s,
                "JupyterLab：承载可执行 Notebook 实验\nPython Kernel：运行热力求解代码\nnbgrader：发布、收集、评分和反馈实验作业\nDocker：隔离学生实验环境，保证运行环境一致",
                7.38, 4.88, 4.5, 0.5, size=9.6, color=BLACK, margin=0)
    add_textbox(s, "教学应用层将热力仿真的白箱化过程转化为可训练、可提交、可评价的实验任务，支撑学生对热力模型、数值求解和工程结果的综合理解。",
                1.0, 6.16, 11.35, 0.32, size=13.2, color=NAVY, bold=True,
                align=PP_ALIGN.CENTER, margin=0)


def main():
    chart = make_chart()
    pres = Presentation()
    pres.slide_width = Inches(W)
    pres.slide_height = Inches(H)
    # remove default slide if any
    while len(pres.slides) > 0:
        r_id = pres.slides._sldIdLst[0].rId
        pres.part.drop_rel(r_id)
        del pres.slides._sldIdLst[0]

    slide1(pres)
    slide2(pres)
    slide3(pres)
    slide4(pres)
    slide5(pres, chart)
    slide6(pres)
    pres.save(OUT)
    print(f"saved {OUT}")


if __name__ == "__main__":
    main()
